from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
import os

from app.tool.base import BaseTool, ToolResult
from app.config import WORKSPACE_ROOT


class PPTCreator(BaseTool):
    """A tool for creating PowerPoint presentations."""

    name: str = "ppt_creator"
    description: str = """Create and edit PowerPoint presentations.
This tool allows you to create slides, add content, and save presentations to files.
You can create new presentations, add various types of slides (title, content, etc.),
and customize text, layouts, and basic formatting.
"""
    parameters: dict = {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": ["create", "add_slide", "add_text", "add_image", "save"],
                "description": "The action to perform on the presentation",
            },
            "file_path": {
                "type": "string",
                "description": "Path where the presentation should be saved or loaded from",
            },
            "slide_type": {
                "type": "string",
                "enum": ["title", "content", "two_content", "section", "blank"],
                "description": "Type of slide to add (used with add_slide action)",
            },
            "title": {
                "type": "string",
                "description": "Title text for the slide (used with add_slide action)",
            },
            "content": {
                "type": "string",
                "description": "Content text or bullet points (used with add_text action)",
            },
            "slide_index": {
                "type": "integer",
                "description": "Index of the slide to modify (0-based)",
            },
            "placeholder_index": {
                "type": "integer",
                "description": "Index of the placeholder to add text to",
            },
            "image_path": {
                "type": "string",
                "description": "Path to the image file to add (used with add_image action)",
            },
        },
        "required": ["action"],
    }

    # 存储当前正在处理的演示文稿
    _current_presentation: Optional[Presentation] = None
    _current_file_path: Optional[str] = None

    # 幻灯片布局索引映射
    _slide_layouts = {
        "title": 0,        # 标题幻灯片
        "content": 1,      # 标题和内容
        "two_content": 3,  # 两栏内容
        "section": 2,      # 章节标题
        "blank": 5,        # 空白幻灯片
    }

    async def execute(
        self,
        action: str,
        file_path: Optional[str] = None,
        slide_type: Optional[str] = None,
        title: Optional[str] = None,
        content: Optional[str] = None,
        slide_index: Optional[int] = None,
        placeholder_index: Optional[int] = None,
        image_path: Optional[str] = None,
        **kwargs
    ) -> ToolResult:
        """Execute the PPT creation tool with the specified parameters."""

        try:
            # 处理文件路径
            if file_path:
                if os.path.isabs(file_path):
                    file_name = os.path.basename(file_path)
                    full_path = os.path.join(WORKSPACE_ROOT, file_name)
                else:
                    full_path = os.path.join(WORKSPACE_ROOT, file_path)

                # 确保目录存在
                os.makedirs(os.path.dirname(full_path), exist_ok=True)

                # 更新当前文件路径
                self._current_file_path = full_path

            # 根据不同操作执行相应功能
            if action == "create":
                return await self._create_presentation(file_path)
            elif action == "add_slide":
                return await self._add_slide(slide_type, title)
            elif action == "add_text":
                return await self._add_text(slide_index, placeholder_index, content)
            elif action == "add_image":
                return await self._add_image(slide_index, image_path)
            elif action == "save":
                return await self._save_presentation()
            else:
                return ToolResult(error=f"Unknown action: {action}")

        except Exception as e:
            return ToolResult(error=f"Error in PPT creation: {str(e)}")

    async def _create_presentation(self, file_path: Optional[str]) -> ToolResult:
        """Create a new presentation."""
        self._current_presentation = Presentation()
        if file_path:
            self._current_file_path = self._current_file_path or file_path

        return ToolResult(output="New presentation created successfully.")

    async def _add_slide(self, slide_type: Optional[str], title: Optional[str]) -> ToolResult:
        """Add a new slide to the presentation."""
        if not self._current_presentation:
            return ToolResult(error="No active presentation. Use 'create' action first.")

        if not slide_type or slide_type not in self._slide_layouts:
            return ToolResult(error=f"Invalid slide type: {slide_type}")

        # 获取布局并添加幻灯片
        layout_idx = self._slide_layouts[slide_type]
        slide_layout = self._current_presentation.slide_layouts[layout_idx]
        slide = self._current_presentation.slides.add_slide(slide_layout)

        # 如果提供了标题且幻灯片有标题占位符，则添加标题
        if title and hasattr(slide, "shapes") and slide.shapes.title:
            slide.shapes.title.text = title

        return ToolResult(output=f"Added new {slide_type} slide with index {len(self._current_presentation.slides) - 1}")

    async def _add_text(
        self, slide_index: Optional[int], placeholder_index: Optional[int], content: Optional[str]
    ) -> ToolResult:
        """Add text to a slide placeholder."""
        if not self._current_presentation:
            return ToolResult(error="No active presentation. Use 'create' action first.")

        if slide_index is None or slide_index < 0 or slide_index >= len(self._current_presentation.slides):
            return ToolResult(error=f"Invalid slide index: {slide_index}")

        if not content:
            return ToolResult(error="No content provided")

        slide = self._current_presentation.slides[slide_index]

        # 如果指定了占位符索引，尝试添加到该占位符
        if placeholder_index is not None:
            try:
                shape = slide.placeholders[placeholder_index]
                shape.text = content
                return ToolResult(output=f"Added text to placeholder {placeholder_index} on slide {slide_index}")
            except (IndexError, KeyError):
                return ToolResult(error=f"Placeholder {placeholder_index} not found on slide {slide_index}")

        # 否则，尝试添加到第一个文本框或内容占位符
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = content
                return ToolResult(output=f"Added text to slide {slide_index}")

        return ToolResult(error=f"No suitable text placeholder found on slide {slide_index}")

    async def _add_image(self, slide_index: Optional[int], image_path: Optional[str]) -> ToolResult:
        """Add an image to a slide."""
        if not self._current_presentation:
            return ToolResult(error="No active presentation. Use 'create' action first.")

        if slide_index is None or slide_index < 0 or slide_index >= len(self._current_presentation.slides):
            return ToolResult(error=f"Invalid slide index: {slide_index}")

        if not image_path:
            return ToolResult(error="No image path provided")

        # 处理图片路径
        if os.path.isabs(image_path):
            full_image_path = image_path
        else:
            full_image_path = os.path.join(WORKSPACE_ROOT, image_path)

        if not os.path.exists(full_image_path):
            return ToolResult(error=f"Image file not found: {full_image_path}")

        slide = self._current_presentation.slides[slide_index]

        # 添加图片到幻灯片
        try:
            slide.shapes.add_picture(
                full_image_path,
                left=Inches(1),
                top=Inches(2),
                width=Inches(4)
            )
            return ToolResult(output=f"Added image to slide {slide_index}")
        except Exception as e:
            return ToolResult(error=f"Error adding image: {str(e)}")

    async def _save_presentation(self) -> ToolResult:
        """Save the presentation to a file."""
        if not self._current_presentation:
            return ToolResult(error="No active presentation. Use 'create' action first.")

        if not self._current_file_path:
            return ToolResult(error="No file path specified for saving")

        # 确保文件扩展名为 .pptx
        if not self._current_file_path.endswith('.pptx'):
            self._current_file_path += '.pptx'

        # 保存演示文稿
        try:
            self._current_presentation.save(self._current_file_path)
            return ToolResult(output=f"Presentation saved to {self._current_file_path}")
        except Exception as e:
            return ToolResult(error=f"Error saving presentation: {str(e)}")
